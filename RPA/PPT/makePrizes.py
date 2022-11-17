import pandas as pd
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN

# 슬라이드에 텍스트 추가하는 함수 정의
def add_text(slide, left, top, width, height, font_size, alignment, text):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    textbox.text_frame.word_wrap = True
    paragraph = textbox.text_frame.paragraphs[0]
    paragraph.font.bold = True
    paragraph.font.size = font_size
    paragraph.alignment = alignment
    paragraph.text = text


# 엑셀 데이터 가져오기
data = pd.read_excel('상장목록.xlsx')
students = data['학생']
titles = data['제목']
contents = data['내용']

# 프레젠테이션 만들기
prs = Presentation()
full_width = prs.slide_width

for i in range(len(data)):
    # 빈 슬라이드 추가
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 상장 배경 이미지 추가
    slide.shapes.add_picture('상장템플릿.png', left=Cm(0), top=Cm(0), width=full_width, height=prs.slide_height)
    # 상장 내용 추가
    add_text(slide, Cm(0), Cm(2), full_width, Cm(0), Pt(50), PP_ALIGN.CENTER, str(titles[i]))
    add_text(slide, Cm(0), Cm(5), full_width - Cm(3), Cm(0), Pt(18), PP_ALIGN.RIGHT, '파이썬 초등학교 3학년 10반')
    add_text(slide, Cm(0), Cm(6), full_width - Cm(3), Cm(0), Pt(30), PP_ALIGN.RIGHT, str(students[i]))
    add_text(slide, Cm(3), Cm(8), full_width - Cm(6), Cm(0), Pt(23), PP_ALIGN.LEFT, str(contents[i]))
    add_text(slide, Cm(0), Cm(13), full_width, Cm(0), Pt(20), PP_ALIGN.CENTER, '2030년 01월 30일')
    add_text(slide, Cm(0), Cm(14.8), full_width, Cm(0), Pt(30), PP_ALIGN.CENTER, '너의 친구 룡룡이')

prs.save("result.pptx")
