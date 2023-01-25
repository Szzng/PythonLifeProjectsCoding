from pptx import Presentation
from pptx.util import Cm, Pt

# 프레젠테이션 만들기
prs = Presentation()

# 슬라이드 레이아웃 종류 확인하기 (0~10)
layouts = prs.slide_layouts
for i in range(len(layouts)):
    print(i, layouts[i].name)

# 첫 번째 슬라이드 (제목 슬라이드)
layout = prs.slide_layouts[0]  # 슬라이드 레이아웃 선택
slide1 = prs.slides.add_slide(layout)  # 선택한 레이아웃으로 슬라이드 추가
shapes = slide1.shapes  # 슬라이드의 도형 가져오기
shapes[0].text = "슬라이드1 제목"  # 첫 번째 도형(제목)에 텍스트 입력
shapes[1].text = "슬라이드1 부제목"  # 두 번째 도형(부제목)에 텍스트 입력

# 두 번째 슬라이드 (제목 + 내용 슬라이드)
slide2 = prs.slides.add_slide(prs.slide_layouts[1])
slide2.shapes[0].text = "슬라이드2 제목"
slide2.shapes[1].text = "슬라이드2 내용"

# 세 번째 슬라이드 (빈 슬라이드 + 이미지 추가)
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
slide3.shapes.add_picture('룡룡.png', left=Cm(9), top=Cm(6))

# 네 번째 슬라이드 (빈 슬라이드 + 텍스트 상자 추가)
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
textbox = slide4.shapes.add_textbox(Cm(2), Cm(2), Cm(15), Cm(0))
textbox.text_frame.word_wrap = True
paragraph = textbox.text_frame.paragraphs[0]
paragraph.font.size = Pt(50)
paragraph.font.bold = True
paragraph.text = "크고 굵은 글씨의 문단을 추가해봅시다."

# 프레젠테이션 파일을 저장하기
prs.save("basic.pptx")
